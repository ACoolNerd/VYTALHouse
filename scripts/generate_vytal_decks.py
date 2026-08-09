import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

def create_pptx_deck(output_filename, title_text, is_investor=False):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    DARK_BG = RGBColor(11, 12, 14)       # #0B0C0E
    GRAPHITE = RGBColor(18, 22, 31)     # #12161F
    SILVER = RGBColor(183, 189, 198)    # #B7BDC6
    LIME = RGBColor(198, 252, 6)        # #C6FC06
    WHITE = RGBColor(255, 255, 255)

    slides_data = [
        {
            "slide_title": "VYTAL HOUSE",
            "subtitle": "RENEW • RESTORE • RECHARGE",
            "content": [
                "Physician-Led Restoration House & Sublingual Performance Ecosystem",
                "AT THE INTERSECTION OF SCIENCE, DISCIPLINE, AND LONGEVITY",
                "Master Brand Direction V1.0 | 2035 Prism Silver Edition"
            ],
            "img": "public/assets/brand/master-os/00-master-overview.jpg"
        },
        {
            "slide_title": "01. EXECUTIVE SUMMARY & MARKET OPPORTUNITY",
            "subtitle": "THE QUANTITATIVE LONGEVITY & SOCIAL WELLNESS MOAT",
            "content": [
                "• TAM/SAM/SOM: $100B+ Longevity & Performance Market Expansion",
                "• The Void: Sober executive social clubs replacing traditional night venues",
                "• Moat: Quantitative Trajectory Loop (Baseline → Protocol → Re-Test → VYTAL Index 0-100)",
                "• Clinical Oversight: 100% Owned Medical PC led by Dr. Abasi Bomani, MD"
            ],
            "img": "public/assets/brand/master-os/06-digital-app-ecosystem.png"
        },
        {
            "slide_title": "02. VYTAL POUCHES PRODUCT ECOSYSTEM",
            "subtitle": "ZERO TOBACCO • ZERO NICOTINE • SUBLINGUAL PRECISION",
            "content": [
                "• VYTAL Energy (Sour Soursop Lime | Citrus Lime #C6FC06 Ring)",
                "• VYTAL Tidal (Caribbean Sea Moss | Blue Raspberry #00F0FF Ring)",
                "• VYTAL NAD+ Support (Sublingual NAD+ | Wild Berry #CDB8FF Ring)",
                "• Packaging: Soft-Touch Matte Black Aluminum Cans & Starter Kit Boxes"
            ],
            "img": "public/assets/brand/master-os/07-product-packaging.png"
        },
        {
            "slide_title": "03. THREE-TIER REVENUE ENGINE & FINANCIAL MODEL",
            "subtitle": "RECURRING MEMBERSHIPS + CLINICAL PC FEES + MSO RETAIL",
            "content": [
                "• Digital Portal ($29/mo) | House Access ($295/mo) | Black VIP ($695/mo)",
                "• Unit Economics: 68% Gross Margin | $4.2M Year 1 Flagship Yield",
                "• Capital Requirement: $2.5M Seed Seed Raise | Payback: 14 Months",
                "• Scalability: 5,000 SF Modular Site Replication Model (MD, VA, NY)"
            ],
            "img": "public/assets/brand/master-os/08-environmental-signage.jpg"
        },
        {
            "slide_title": "04. BRAND GOVERNANCE & LOGO STANDARDS",
            "subtitle": "11 OFFICIAL LOGO CONFIGURATIONS & MATERIAL PALETTE",
            "content": [
                "• Master Monogram: 3D Iridescent Platinum Silver V Mark",
                "• Primary Color System: Obsidian (#0B0C0E), Graphite (#12161F), Titanium (#B7BDC6)",
                "• Signal Discipline: Acid Lime (#C6FC06) strictly for CTAs & Data Signals",
                "• Production Standard: Zero WIP Labels | 100% Master Brand OS Compliant"
            ],
            "img": "public/assets/brand/master-os/02-logo-standards.png"
        }
    ]

    for data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

        # Title
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = data["slide_title"]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        p2 = tf.add_paragraph()
        p2.text = data["subtitle"]
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = LIME

        # Content Box
        tx_box2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.0), Inches(4.5))
        tf2 = tx_box2.text_frame
        tf2.word_wrap = True
        for idx, line in enumerate(data["content"]):
            p_c = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
            p_c.text = line
            p_c.font.size = Pt(14)
            p_c.font.color.rgb = SILVER

        # Image
        if os.path.exists(data["img"]):
            slide.shapes.add_picture(data["img"], Inches(7.2), Inches(2.0), width=Inches(5.3))

    prs.save(output_filename)
    print(f"Generated PPTX: {output_filename}")

def create_pdf_deck(output_filename, title_text):
    doc = SimpleDocTemplate(
        output_filename,
        pagesize=landscape(letter),
        rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36
    )
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=24,
        textColor=colors.HexColor('#FFFFFF'),
        spaceAfter=6
    )
    sub_style = ParagraphStyle(
        'DocSub',
        fontName='Helvetica-Bold',
        fontSize=12,
        textColor=colors.HexColor('#C6FC06'),
        spaceAfter=15
    )
    body_style = ParagraphStyle(
        'BodyDark',
        fontName='Helvetica',
        fontSize=11,
        textColor=colors.HexColor('#B7BDC6'),
        spaceBefore=4, spaceAfter=4
    )

    story = []

    pages = [
        ("VYTAL HOUSE — MASTER BRAND PRESENTATION", "RENEW • RESTORE • RECHARGE", [
            "Physician-Led Restoration House & Performance Ecosystem",
            "AT THE INTERSECTION OF SCIENCE, DISCIPLINE, AND LONGEVITY",
            "Master Brand Direction V1.0 | 2035 Prism Silver Edition",
            "Scorecard Ranking: 9.9/10 Institutional Investment Grade"
        ], "public/assets/brand/master-os/00-master-overview.jpg"),

        ("01. EXECUTIVE SUMMARY & MARKET OPPORTUNITY", "THE QUANTITATIVE LONGEVITY & SOCIAL WELLNESS MOAT", [
            "• TAM/SAM/SOM: $100B+ Longevity & Performance Market Expansion",
            "• The Void: Sober executive social clubs replacing traditional night venues",
            "• Moat: Quantitative Trajectory Loop (Baseline → Protocol → Re-Test → VYTAL Index 0-100)",
            "• Clinical Oversight: 100% Owned Medical PC led by Dr. Abasi Bomani, MD"
        ], "public/assets/brand/master-os/06-digital-app-ecosystem.png"),

        ("02. VYTAL POUCHES PRODUCT ECOSYSTEM", "ZERO TOBACCO • ZERO NICOTINE • SUBLINGUAL PRECISION", [
            "• VYTAL Energy (Sour Soursop Lime | Citrus Lime #C6FC06 Ring)",
            "• VYTAL Tidal (Caribbean Sea Moss | Blue Raspberry #00F0FF Ring)",
            "• VYTAL NAD+ Support (Sublingual NAD+ | Wild Berry #CDB8FF Ring)",
            "• Packaging: Soft-Touch Matte Black Aluminum Cans & Starter Kit Boxes"
        ], "public/assets/brand/master-os/07-product-packaging.png"),

        ("03. THREE-TIER REVENUE ENGINE & FINANCIAL MODEL", "RECURRING MEMBERSHIPS + CLINICAL PC FEES + MSO RETAIL", [
            "• Digital Portal ($29/mo) | House Access ($295/mo) | Black VIP ($695/mo)",
            "• Unit Economics: 68% Gross Margin | $4.2M Year 1 Flagship Yield",
            "• Capital Requirement: $2.5M Seed Raise | Payback: 14 Months",
            "• Scalability: 5,000 SF Modular Site Replication Model (MD, VA, NY)"
        ], "public/assets/brand/master-os/08-environmental-signage.jpg"),

        ("04. BRAND GOVERNANCE & LOGO STANDARDS", "11 OFFICIAL LOGO CONFIGURATIONS & MATERIAL PALETTE", [
            "• Master Monogram: 3D Iridescent Platinum Silver V Mark",
            "• Primary Color System: Obsidian (#0B0C0E), Graphite (#12161F), Titanium (#B7BDC6)",
            "• Signal Discipline: Acid Lime (#C6FC06) strictly for CTAs & Data Signals",
            "• Production Standard: Zero WIP Labels | 100% Master Brand OS Compliant"
        ], "public/assets/brand/master-os/02-logo-standards.png")
    ]

    for idx, (t, sub, bullets, img_path) in enumerate(pages):
        story.append(Paragraph(t, title_style))
        story.append(Paragraph(sub, sub_style))
        
        bullet_pars = [Paragraph(b, body_style) for b in bullets]
        
        if os.path.exists(img_path):
            img = Image(img_path, width=360, height=200)
            t_data = [[bullet_pars, img]]
            table = Table(t_data, colWidths=[360, 360])
            table.setStyle(TableStyle([
                ('VALIGN', (0,0), (-1,-1), 'TOP'),
                ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#0B0C0E')),
                ('TEXTCOLOR', (0,0), (-1,-1), colors.HexColor('#B7BDC6')),
            ]))
            story.append(table)
        else:
            for b in bullet_pars:
                story.append(b)
                
        if idx < len(pages) - 1:
            story.append(PageBreak())

    def draw_bg(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(colors.HexColor('#0B0C0E'))
        canvas.rect(0, 0, doc.pagesize[0], doc.pagesize[1], fill=True, stroke=False)
        canvas.restoreState()

    doc.build(story, onFirstPage=draw_bg, onLaterPages=draw_bg)
    print(f"Generated PDF: {output_filename}")

if __name__ == "__main__":
    os.makedirs("dist/presentations", exist_ok=True)
    
    # 1. Generate PPTX decks
    create_pptx_deck("dist/presentations/VYTAL_Pouches_Faces_Brand_Pitch.pptx", "VYTAL POUCHES BRAND PITCH")
    create_pptx_deck("dist/presentations/VYTAL_Pouches_Investor_Deck.pptx", "VYTAL HOUSE INVESTOR DECK", is_investor=True)
    
    # 2. Generate PDF decks
    create_pdf_deck("dist/presentations/VYTAL_Pouches_Faces_Brand_Pitch.pdf", "VYTAL POUCHES BRAND PITCH")
    create_pdf_deck("dist/presentations/VYTAL_Pouches_Investor_Deck.pdf", "VYTAL HOUSE INVESTOR DECK")
